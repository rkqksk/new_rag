"""
문서 수집 및 처리 서비스
- 다양한 형식 지원 (PDF, Excel, Images, Text, HTML, XML, CSV, JSON, 등)
- 자동 청킹 및 메타데이터 추출
- Qdrant 벡터 DB 저장소 연동
"""

import asyncio
import csv
import hashlib
import json
import logging
import os
import xml.etree.ElementTree as ET
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional
from uuid import uuid4

import aiohttp
import pandas as pd
import pytesseract
import requests
from bs4 import BeautifulSoup
from PIL import Image
from qdrant_client.models import Distance, PointStruct, VectorParams
from sentence_transformers import SentenceTransformer
from unstructured.partition.html import partition_html
from unstructured.partition.image import partition_image
from unstructured.partition.pdf import partition_pdf

logger = logging.getLogger(__name__)


class DocumentChunk:
    """문서 청크 데이터 모델"""

    def __init__(self, text: str, doc_id: str, chunk_index: int, metadata: Dict[str, Any] = None):
        self.id = str(uuid4())
        self.text = text
        self.doc_id = doc_id
        self.chunk_index = chunk_index
        self.metadata = metadata or {}
        self.created_at = datetime.utcnow().isoformat()

    def to_dict(self):
        return {
            "id": self.id,
            "text": self.text,
            "doc_id": self.doc_id,
            "chunk_index": self.chunk_index,
            "metadata": self.metadata,
            "created_at": self.created_at,
        }


class DocumentIngestionService:
    """
    다양한 형식의 문서를 처리하고 Qdrant에 저장하는 서비스
    """

    def __init__(
        self, qdrant_client, embedding_model: str = "sentence-transformers/all-MiniLM-L6-v2"
    ):
        self.qdrant_client = qdrant_client
        self.embedding_model = SentenceTransformer(embedding_model)
        self.collection_name = "documents"
        self.chunk_size = 512  # 토큰 근사치
        self.chunk_overlap = 50
        self.vector_size = self.embedding_model.get_sentence_embedding_dimension()

        # 컬렉션 초기화
        self._init_collection()

    def _init_collection(self):
        """Qdrant 컬렉션 초기화"""
        try:
            self.qdrant_client.get_collection(self.collection_name)
        except Exception:
            # 컬렉션이 없으면 생성
            self.qdrant_client.create_collection(
                collection_name=self.collection_name,
                vectors_config=VectorParams(size=self.vector_size, distance=Distance.COSINE),
            )
            logger.info(f"Created Qdrant collection: {self.collection_name}")

    async def ingest_file(
        self, file_path: str, doc_type: Optional[str] = None, metadata: Dict[str, Any] = None
    ) -> Dict[str, Any]:
        """
        파일을 업로드하고 처리

        Args:
            file_path: 파일 경로
            doc_type: 문서 타입 (auto-detect if None)
            metadata: 추가 메타데이터

        Returns:
            처리 결과 (doc_id, chunks_count, vectors_stored)
        """
        metadata = metadata or {}
        doc_id = str(uuid4())

        try:
            # 파일 타입 자동 감지
            if doc_type is None:
                doc_type = self._detect_file_type(file_path)

            logger.info(f"Ingesting {doc_type}: {file_path}")

            # 문서 타입별 처리
            if doc_type == "pdf":
                chunks = await self._process_pdf(file_path, doc_id, metadata)
            elif doc_type == "excel":
                chunks = await self._process_excel(file_path, doc_id, metadata)
            elif doc_type == "csv":
                chunks = await self._process_csv(file_path, doc_id, metadata)
            elif doc_type == "xml":
                chunks = await self._process_xml(file_path, doc_id, metadata)
            elif doc_type == "json":
                chunks = await self._process_json(file_path, doc_id, metadata)
            elif doc_type == "image":
                chunks = await self._process_image(file_path, doc_id, metadata)
            elif doc_type == "html":
                chunks = await self._process_html(file_path, doc_id, metadata)
            elif doc_type == "word":
                chunks = await self._process_word(file_path, doc_id, metadata)
            elif doc_type == "powerpoint":
                chunks = await self._process_powerpoint(file_path, doc_id, metadata)
            elif doc_type == "archive":
                chunks = await self._process_archive(file_path, doc_id, metadata)
            elif doc_type == "text":
                chunks = await self._process_text(file_path, doc_id, metadata)
            else:
                chunks = await self._process_text(file_path, doc_id, metadata)

            # 벡터 생성 및 저장
            stored_count = await self._store_vectors(chunks, doc_id)

            result = {
                "doc_id": doc_id,
                "doc_type": doc_type,
                "file_name": Path(file_path).name,
                "chunks_count": len(chunks),
                "vectors_stored": stored_count,
                "metadata": metadata,
                "processed_at": datetime.utcnow().isoformat(),
            }

            logger.info(f"Successfully ingested {len(chunks)} chunks from {file_path}")
            return result

        except Exception as e:
            logger.error(f"Error ingesting file {file_path}: {e}")
            raise

    def _detect_file_type(self, file_path: str) -> str:
        """파일 확장자로 타입 감지"""
        ext = Path(file_path).suffix.lower()

        pdf_exts = {".pdf"}
        excel_exts = {".xlsx", ".xls"}
        csv_exts = {".csv"}
        xml_exts = {".xml"}
        json_exts = {".json"}
        image_exts = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".webp"}
        html_exts = {".html", ".htm"}
        text_exts = {".txt", ".md", ".rst", ".log"}
        word_exts = {".docx", ".doc"}
        powerpoint_exts = {".pptx", ".ppt"}
        archive_exts = {".zip", ".tar", ".gz"}

        if ext in pdf_exts:
            return "pdf"
        elif ext in excel_exts:
            return "excel"
        elif ext in csv_exts:
            return "csv"
        elif ext in xml_exts:
            return "xml"
        elif ext in json_exts:
            return "json"
        elif ext in image_exts:
            return "image"
        elif ext in html_exts:
            return "html"
        elif ext in text_exts:
            return "text"
        elif ext in word_exts:
            return "word"
        elif ext in powerpoint_exts:
            return "powerpoint"
        elif ext in archive_exts:
            return "archive"
        else:
            return "text"  # 기본값

    async def _process_pdf(
        self, file_path: str, doc_id: str, metadata: Dict[str, Any]
    ) -> List[DocumentChunk]:
        """PDF 파일 처리"""
        chunks = []
        metadata.update(
            {
                "source": "pdf",
                "file_path": file_path,
            }
        )

        try:
            # unstructured를 사용한 PDF 파싱
            elements = partition_pdf(file_path)

            text_buffer = []
            token_count = 0

            for elem in elements:
                text = elem.text.strip()
                if not text:
                    continue

                # 토큰 추정 (단어 수 * 1.3)
                token_estimate = len(text.split()) * 1.3

                if token_count + token_estimate > self.chunk_size and text_buffer:
                    # 청크 생성
                    chunk_text = " ".join(text_buffer)
                    chunk = DocumentChunk(
                        text=chunk_text,
                        doc_id=doc_id,
                        chunk_index=len(chunks),
                        metadata=metadata.copy(),
                    )
                    chunks.append(chunk)
                    text_buffer = []
                    token_count = 0

                text_buffer.append(text)
                token_count += token_estimate

            # 남은 텍스트 처리
            if text_buffer:
                chunk_text = " ".join(text_buffer)
                chunk = DocumentChunk(
                    text=chunk_text,
                    doc_id=doc_id,
                    chunk_index=len(chunks),
                    metadata=metadata.copy(),
                )
                chunks.append(chunk)

        except Exception as e:
            logger.error(f"Error processing PDF {file_path}: {e}")
            raise

        return chunks

    async def _process_excel(
        self, file_path: str, doc_id: str, metadata: Dict[str, Any]
    ) -> List[DocumentChunk]:
        """Excel 파일 처리"""
        chunks = []
        metadata.update(
            {
                "source": "excel",
                "file_path": file_path,
            }
        )

        try:
            # CSV도 같은 방식으로 처리
            if file_path.endswith(".csv"):
                df = pd.read_csv(file_path)
            else:
                df = pd.read_excel(file_path)

            # 각 행을 텍스트로 변환
            for sheet_name in ([None] if isinstance(df, pd.DataFrame) else df.sheet_names):
                if sheet_name is not None and file_path.endswith((".xlsx", ".xls")):
                    df = pd.read_excel(file_path, sheet_name=sheet_name)

                for idx, row in df.iterrows():
                    # 행의 모든 열을 텍스트로 변환
                    row_text = " | ".join(f"{col}: {val}" for col, val in row.items())

                    chunk = DocumentChunk(
                        text=row_text,
                        doc_id=doc_id,
                        chunk_index=len(chunks),
                        metadata={**metadata, "row_index": idx, "sheet_name": sheet_name},
                    )
                    chunks.append(chunk)

        except Exception as e:
            logger.error(f"Error processing Excel {file_path}: {e}")
            raise

        return chunks

    async def _process_image(
        self, file_path: str, doc_id: str, metadata: Dict[str, Any]
    ) -> List[DocumentChunk]:
        """이미지 파일 처리 (OCR)"""
        chunks = []
        metadata.update(
            {
                "source": "image",
                "file_path": file_path,
            }
        )

        try:
            # unstructured의 이미지 파싱
            elements = partition_image(file_path)

            text_parts = []
            for elem in elements:
                text = elem.text.strip()
                if text:
                    text_parts.append(text)

            if text_parts:
                combined_text = " ".join(text_parts)
                chunk = DocumentChunk(
                    text=combined_text, doc_id=doc_id, chunk_index=0, metadata=metadata.copy()
                )
                chunks.append(chunk)

        except Exception as e:
            logger.error(f"Error processing image {file_path}: {e}")
            raise

        return chunks

    async def _process_html(
        self, file_path: str, doc_id: str, metadata: Dict[str, Any]
    ) -> List[DocumentChunk]:
        """HTML 파일 처리"""
        chunks = []
        metadata.update(
            {
                "source": "html",
                "file_path": file_path,
            }
        )

        try:
            # 파일에서 읽거나 URL에서 가져오기
            if file_path.startswith("http"):
                async with aiohttp.ClientSession() as session:
                    async with session.get(file_path) as resp:
                        html_content = await resp.text()
            else:
                with open(file_path, "r", encoding="utf-8") as f:
                    html_content = f.read()

            # unstructured를 사용한 HTML 파싱
            elements = partition_html(text=html_content)

            text_buffer = []
            token_count = 0

            for elem in elements:
                text = elem.text.strip()
                if not text:
                    continue

                token_estimate = len(text.split()) * 1.3

                if token_count + token_estimate > self.chunk_size and text_buffer:
                    chunk_text = " ".join(text_buffer)
                    chunk = DocumentChunk(
                        text=chunk_text,
                        doc_id=doc_id,
                        chunk_index=len(chunks),
                        metadata=metadata.copy(),
                    )
                    chunks.append(chunk)
                    text_buffer = []
                    token_count = 0

                text_buffer.append(text)
                token_count += token_estimate

            if text_buffer:
                chunk_text = " ".join(text_buffer)
                chunk = DocumentChunk(
                    text=chunk_text,
                    doc_id=doc_id,
                    chunk_index=len(chunks),
                    metadata=metadata.copy(),
                )
                chunks.append(chunk)

        except Exception as e:
            logger.error(f"Error processing HTML {file_path}: {e}")
            raise

        return chunks

    async def _process_text(
        self, file_path: str, doc_id: str, metadata: Dict[str, Any]
    ) -> List[DocumentChunk]:
        """텍스트 파일 처리"""
        chunks = []
        metadata.update(
            {
                "source": "text",
                "file_path": file_path,
            }
        )

        try:
            with open(file_path, "r", encoding="utf-8") as f:
                text = f.read()

            # 고정 크기로 청킹
            words = text.split()
            chunk_words = int(self.chunk_size / 1.3)  # 토큰을 단어로 변환

            for i in range(0, len(words), chunk_words - int(self.chunk_overlap / 1.3)):
                chunk_text = " ".join(words[i : i + chunk_words])
                if chunk_text.strip():
                    chunk = DocumentChunk(
                        text=chunk_text,
                        doc_id=doc_id,
                        chunk_index=len(chunks),
                        metadata=metadata.copy(),
                    )
                    chunks.append(chunk)

        except Exception as e:
            logger.error(f"Error processing text file {file_path}: {e}")
            raise

        return chunks

    async def _process_csv(
        self, file_path: str, doc_id: str, metadata: Dict[str, Any]
    ) -> List[DocumentChunk]:
        """CSV 파일 처리"""
        chunks = []
        metadata.update(
            {
                "source": "csv",
                "file_path": file_path,
            }
        )

        try:
            df = pd.read_csv(file_path)

            # 헤더를 첫 청크에 포함
            headers = " | ".join(df.columns.tolist())
            header_chunk = DocumentChunk(
                text=f"CSV Headers: {headers}",
                doc_id=doc_id,
                chunk_index=0,
                metadata={**metadata, "chunk_type": "headers"},
            )
            chunks.append(header_chunk)

            # 각 행을 청크로 변환
            for idx, row in df.iterrows():
                row_text = " | ".join(f"{col}: {str(val)}" for col, val in row.items())

                chunk = DocumentChunk(
                    text=row_text,
                    doc_id=doc_id,
                    chunk_index=len(chunks),
                    metadata={**metadata, "row_index": int(idx), "chunk_type": "data_row"},
                )
                chunks.append(chunk)

        except Exception as e:
            logger.error(f"Error processing CSV {file_path}: {e}")
            raise

        return chunks

    async def _process_xml(
        self, file_path: str, doc_id: str, metadata: Dict[str, Any]
    ) -> List[DocumentChunk]:
        """XML 파일 처리"""
        chunks = []
        metadata.update(
            {
                "source": "xml",
                "file_path": file_path,
            }
        )

        try:
            tree = ET.parse(file_path)
            root = tree.getroot()

            # XML을 텍스트로 변환
            def xml_to_text(element, depth=0):
                texts = []
                if element.text and element.text.strip():
                    texts.append("  " * depth + element.text.strip())
                for child in element:
                    texts.extend(xml_to_text(child, depth + 1))
                if element.tail and element.tail.strip():
                    texts.append("  " * depth + element.tail.strip())
                return texts

            xml_texts = xml_to_text(root)
            text_buffer = []
            token_count = 0

            for text in xml_texts:
                token_estimate = len(text.split()) * 1.3

                if token_count + token_estimate > self.chunk_size and text_buffer:
                    chunk_text = "\n".join(text_buffer)
                    chunk = DocumentChunk(
                        text=chunk_text,
                        doc_id=doc_id,
                        chunk_index=len(chunks),
                        metadata=metadata.copy(),
                    )
                    chunks.append(chunk)
                    text_buffer = []
                    token_count = 0

                text_buffer.append(text)
                token_count += token_estimate

            if text_buffer:
                chunk_text = "\n".join(text_buffer)
                chunk = DocumentChunk(
                    text=chunk_text,
                    doc_id=doc_id,
                    chunk_index=len(chunks),
                    metadata=metadata.copy(),
                )
                chunks.append(chunk)

        except Exception as e:
            logger.error(f"Error processing XML {file_path}: {e}")
            raise

        return chunks

    async def _process_json(
        self, file_path: str, doc_id: str, metadata: Dict[str, Any]
    ) -> List[DocumentChunk]:
        """JSON 파일 처리"""
        chunks = []
        metadata.update(
            {
                "source": "json",
                "file_path": file_path,
            }
        )

        try:
            with open(file_path, "r", encoding="utf-8") as f:
                data = json.load(f)

            # JSON을 포맷된 텍스트로 변환
            json_text = json.dumps(data, indent=2, ensure_ascii=False)
            text_buffer = []
            token_count = 0

            for line in json_text.split("\n"):
                if not line.strip():
                    continue

                token_estimate = len(line.split()) * 1.3

                if token_count + token_estimate > self.chunk_size and text_buffer:
                    chunk_text = "\n".join(text_buffer)
                    chunk = DocumentChunk(
                        text=chunk_text,
                        doc_id=doc_id,
                        chunk_index=len(chunks),
                        metadata=metadata.copy(),
                    )
                    chunks.append(chunk)
                    text_buffer = []
                    token_count = 0

                text_buffer.append(line)
                token_count += token_estimate

            if text_buffer:
                chunk_text = "\n".join(text_buffer)
                chunk = DocumentChunk(
                    text=chunk_text,
                    doc_id=doc_id,
                    chunk_index=len(chunks),
                    metadata=metadata.copy(),
                )
                chunks.append(chunk)

        except Exception as e:
            logger.error(f"Error processing JSON {file_path}: {e}")
            raise

        return chunks

    async def _process_word(
        self, file_path: str, doc_id: str, metadata: Dict[str, Any]
    ) -> List[DocumentChunk]:
        """Word 파일 처리 (.docx, .doc)"""
        chunks = []
        metadata.update(
            {
                "source": "word",
                "file_path": file_path,
            }
        )

        try:
            from docx import Document as DocxDocument

            doc = DocxDocument(file_path)
            text_buffer = []
            token_count = 0

            for para in doc.paragraphs:
                text = para.text.strip()
                if not text:
                    continue

                token_estimate = len(text.split()) * 1.3

                if token_count + token_estimate > self.chunk_size and text_buffer:
                    chunk_text = " ".join(text_buffer)
                    chunk = DocumentChunk(
                        text=chunk_text,
                        doc_id=doc_id,
                        chunk_index=len(chunks),
                        metadata=metadata.copy(),
                    )
                    chunks.append(chunk)
                    text_buffer = []
                    token_count = 0

                text_buffer.append(text)
                token_count += token_estimate

            if text_buffer:
                chunk_text = " ".join(text_buffer)
                chunk = DocumentChunk(
                    text=chunk_text,
                    doc_id=doc_id,
                    chunk_index=len(chunks),
                    metadata=metadata.copy(),
                )
                chunks.append(chunk)

        except ImportError:
            logger.warning("python-docx not installed, treating as text file")
            chunks = await self._process_text(file_path, doc_id, metadata)
        except Exception as e:
            logger.error(f"Error processing Word {file_path}: {e}")
            raise

        return chunks

    async def _process_powerpoint(
        self, file_path: str, doc_id: str, metadata: Dict[str, Any]
    ) -> List[DocumentChunk]:
        """PowerPoint 파일 처리 (.pptx, .ppt)"""
        chunks = []
        metadata.update(
            {
                "source": "powerpoint",
                "file_path": file_path,
            }
        )

        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            slide_num = 0

            for slide in prs.slides:
                slide_num += 1
                slide_text = f"=== Slide {slide_num} ===\n"

                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + "\n"

                if slide_text.strip():
                    chunk = DocumentChunk(
                        text=slide_text.strip(),
                        doc_id=doc_id,
                        chunk_index=len(chunks),
                        metadata={**metadata, "slide_number": slide_num},
                    )
                    chunks.append(chunk)

        except ImportError:
            logger.warning("python-pptx not installed, treating as text file")
            chunks = await self._process_text(file_path, doc_id, metadata)
        except Exception as e:
            logger.error(f"Error processing PowerPoint {file_path}: {e}")
            raise

        return chunks

    async def _process_archive(
        self, file_path: str, doc_id: str, metadata: Dict[str, Any]
    ) -> List[DocumentChunk]:
        """압축 파일 처리 (.zip, .tar, .gz)"""
        chunks = []
        metadata.update(
            {
                "source": "archive",
                "file_path": file_path,
            }
        )

        try:
            import tarfile
            import tempfile
            import zipfile

            temp_extract_dir = tempfile.mkdtemp()

            try:
                # ZIP 파일 처리
                if file_path.endswith(".zip"):
                    with zipfile.ZipFile(file_path, "r") as zip_ref:
                        zip_ref.extractall(temp_extract_dir)

                # TAR 파일 처리
                elif file_path.endswith((".tar", ".tar.gz", ".tgz")):
                    with tarfile.open(file_path, "r:*") as tar_ref:
                        tar_ref.extractall(temp_extract_dir)

                # 추출된 파일들 처리
                for root, dirs, files in os.walk(temp_extract_dir):
                    for file in files:
                        file_path_inner = os.path.join(root, file)
                        try:
                            file_type = self._detect_file_type(file_path_inner)

                            if file_type == "pdf":
                                inner_chunks = await self._process_pdf(
                                    file_path_inner, doc_id, metadata
                                )
                            elif file_type == "text":
                                inner_chunks = await self._process_text(
                                    file_path_inner, doc_id, metadata
                                )
                            elif file_type == "csv":
                                inner_chunks = await self._process_csv(
                                    file_path_inner, doc_id, metadata
                                )
                            elif file_type == "json":
                                inner_chunks = await self._process_json(
                                    file_path_inner, doc_id, metadata
                                )
                            else:
                                inner_chunks = await self._process_text(
                                    file_path_inner, doc_id, metadata
                                )

                            chunks.extend(inner_chunks)
                        except Exception as e:
                            logger.warning(f"Failed to process {file_path_inner}: {e}")

            finally:
                # 임시 디렉토리 정리
                if os.path.exists(temp_extract_dir):
                    import shutil

                    shutil.rmtree(temp_extract_dir)

        except Exception as e:
            logger.error(f"Error processing archive {file_path}: {e}")
            raise

        return chunks

    async def _store_vectors(self, chunks: List[DocumentChunk], doc_id: str) -> int:
        """
        청크들을 임베딩하고 Qdrant에 저장

        Returns:
            저장된 벡터 개수
        """
        if not chunks:
            return 0

        try:
            # 배치 임베딩
            texts = [chunk.text for chunk in chunks]
            embeddings = self.embedding_model.encode(texts, show_progress_bar=False)

            # Qdrant Point 생성
            points = []
            for i, (chunk, embedding) in enumerate(zip(chunks, embeddings)):
                # Point ID: UUID의 해시로 양수 정수 생성
                point_id = hash(chunk.id) & 0x7FFFFFFF

                points.append(
                    PointStruct(
                        id=point_id,
                        vector=embedding.tolist(),
                        payload={
                            "chunk_id": chunk.id,
                            "doc_id": chunk.doc_id,
                            "text": chunk.text,
                            "chunk_index": chunk.chunk_index,
                            "metadata": chunk.metadata,
                            "created_at": chunk.created_at,
                        },
                    )
                )

            # Qdrant에 저장
            self.qdrant_client.upsert(collection_name=self.collection_name, points=points)

            logger.info(f"Stored {len(points)} vectors for doc_id {doc_id}")
            return len(points)

        except Exception as e:
            logger.error(f"Error storing vectors: {e}")
            raise

    async def search_documents(self, query: str, top_k: int = 5) -> List[Dict[str, Any]]:
        """
        문서 검색

        Args:
            query: 검색 쿼리
            top_k: 반환할 상위 결과 개수

        Returns:
            검색 결과 리스트
        """
        try:
            # 쿼리 임베딩
            query_embedding = self.embedding_model.encode(query).tolist()

            # Qdrant 검색
            search_results = self.qdrant_client.search(
                collection_name=self.collection_name,
                query_vector=query_embedding,
                limit=top_k,
                with_payload=True,
            )

            results = []
            for point in search_results:
                results.append(
                    {
                        "score": point.score,
                        "chunk_id": point.payload.get("chunk_id"),
                        "doc_id": point.payload.get("doc_id"),
                        "text": point.payload.get("text"),
                        "metadata": point.payload.get("metadata"),
                        "created_at": point.payload.get("created_at"),
                    }
                )

            return results

        except Exception as e:
            logger.error(f"Error searching documents: {e}")
            raise

    def get_collection_stats(self) -> Dict[str, Any]:
        """컬렉션 통계 조회"""
        try:
            collection_info = self.qdrant_client.get_collection(self.collection_name)

            return {
                "collection_name": self.collection_name,
                "points_count": collection_info.points_count,
                "vectors_count": collection_info.vectors_count,
                "vector_size": self.vector_size,
                "indexed": (
                    collection_info.indexed_vectors_count
                    if hasattr(collection_info, "indexed_vectors_count")
                    else None
                ),
            }
        except Exception as e:
            logger.error(f"Error getting collection stats: {e}")
            raise
